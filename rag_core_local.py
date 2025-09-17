# rag_core_local.py — RAG sans Elasticsearch (SQLite FTS5 OU Whoosh)

import os, json, re, hashlib, sqlite3
from pathlib import Path
from typing import List, Dict, Optional
import pdfplumber, docx
from pptx import Presentation
import pandas as pd

# --------- Config ----------
RAG_MODE   = os.getenv("RAG_MODE", "sqlite")          # "sqlite" | "whoosh"
DATA_ROOT  = Path(os.getenv("RAG_DATA", "/app/docs")) # dossier à indexer (si on indexe ici)
CACHE_PATH = Path("index_cache.json")

# SQLite
RAG_DB_PATH = os.getenv("RAG_DB_PATH", "rag.db")

# Whoosh
WHOOSH_INDEX_DIR = Path(os.getenv("WHOOSH_INDEX_DIR", "whoosh_index"))

# --------- Extraction texte ----------
SKIP_DIRS = {".git",".terraform","node_modules","__pycache__",".venv",".idea",".vscode"}
TEXT_EXT  = {".txt",".md",".log",".json",".yml",".yaml",".py",".csv"}
DOC_EXT   = {".pdf",".docx",".pptx",".xlsx",".xls"}

def read_text(p: Path) -> str:
    try: return p.read_text(encoding="utf-8", errors="ignore")
    except Exception: return p.read_bytes().decode("latin-1", errors="ignore")

def read_pdf(p: Path) -> str:
    out=[]; 
    with pdfplumber.open(p) as pdf:
        for page in pdf.pages: out.append(page.extract_text() or "")
    return "\n".join(out)

def read_docx(p: Path) -> str:
    d = docx.Document(str(p))
    return "\n".join(par.text for par in d.paragraphs)

def read_pptx(p: Path) -> str:
    prs = Presentation(str(p)); t=[]
    for sl in prs.slides:
        for sh in sl.shapes:
            if hasattr(sh, "text"): t.append(sh.text)
    return "\n".join(t)

def read_xl(p: Path, max_rows=200, max_cols=40) -> str:
    try:
        df = pd.read_csv(p, nrows=max_rows) if p.suffix.lower()==".csv" else pd.read_excel(p, nrows=max_rows)
        df = df.iloc[:, :max_cols]
        return df.to_csv(index=False)
    except Exception:
        return ""

def extract_text(p: Path) -> str:
    ext = p.suffix.lower()
    if ext==".pdf": return read_pdf(p)
    if ext==".docx": return read_docx(p)
    if ext==".pptx": return read_pptx(p)
    if ext in {".xlsx",".xls",".csv"}: return read_xl(p)
    return read_text(p)

def sig_for(p: Path) -> str:
    st = p.stat()
    return f"{int(st.st_mtime)}-{st.st_size}"

def load_cache() -> Dict[str, str]:
    return json.loads(CACHE_PATH.read_text(encoding="utf-8")) if CACHE_PATH.exists() else {}
def save_cache(d: Dict[str, str]):
    CACHE_PATH.write_text(json.dumps(d, indent=2, ensure_ascii=False), encoding="utf-8")

# --------- Rerank sémantique ----------
from sentence_transformers import SentenceTransformer
import numpy as np
from sklearn.metrics.pairwise import cosine_similarity

class SmartVectorizer:
    def __init__(self, model_name="all-MiniLM-L6-v2"):
        self.m = SentenceTransformer(model_name); self.cache={}
    def _enc(self, t:str):
        t=(t or "").strip()
        if t in self.cache: return self.cache[t]
        v=self.m.encode(t, normalize_embeddings=True); self.cache[t]=v; return v
    def rerank(self, q:str, hits:List[Dict], top_n:int=6)->List[Dict]:
        if not hits: return []
        qv=self._enc(q)
        def clip(x,n=1200): return (x or "")[:n]
        X=[self._enc(clip(h["_source"].get("content",""))) for h in hits]
        sims=cosine_similarity(qv.reshape(1,-1), np.vstack(X))[0]
        for h,s in zip(hits, sims): h["_rerank"]=float(s)
        hits.sort(key=lambda h:(h.get("_rerank",0.0), h.get("_score",0.0)), reverse=True)
        return hits[:top_n]

# --------- LLM (OpenAI) avec citations ---------
def build_context(hits: List[Dict], max_chars=8000):
    chunks, meta, total = [], [], 0
    for i,h in enumerate(hits,1):
        src=h["_source"]; txt=(src.get("content") or "").strip().replace("\n"," ")
        if not txt: continue
        block=f"[S{i}] {txt[:1500]}"
        if total+len(block)>max_chars: break
        chunks.append(block); total+=len(block)
        meta.append({"id":f"S{i}","title":src.get("filename",""),"path":src.get("path","")})
    return "\n\n".join(chunks), meta

def format_sources(items: List[Dict]) -> str:
    out=[]
    for s in items:
        if "id" in s:
            out.append(f"- **[{s['id']}]** {s.get('title','')} — {s.get('path','')}")
        else:
            src=s["_source"]; out.append(f"- {src.get('filename','')} — {src.get('path','')}")
    return "\n".join(out)

def answer_with_openai(question: str, hits: List[Dict], model=None) -> str:
    if not os.getenv("OPENAI_API_KEY"):
        return "🔒 OPENAI_API_KEY non défini — réponse LLM désactivée.\n" + format_sources(hits)
    from openai import OpenAI
    client = OpenAI()
    ctx, meta = build_context(hits)
    system = ("Tu es un assistant RAG francophone. Réponds brièvement et précisément, "
              "en citant avec [S1], [S2], etc., et n'utilise que le contexte fourni.")
    model = model or os.getenv("OPENAI_MODEL", "gpt-4o-mini")
    resp = client.chat.completions.create(
        model=model,
        messages=[{"role":"system","content":system},
                  {"role":"user","content":f"Question: {question}\n\nContexte:\n{ctx}\n\nConsignes: réponds en français, ajoute des [Sx], dis si le contexte est insuffisant."}],
        max_tokens=800, temperature=0.3
    )
    txt = resp.choices[0].message.content
    return txt + "\n\n**Sources**\n" + format_sources(meta)

# =========================================================
#                BACKEND A) SQLite FTS5
# =========================================================
SQL_SCHEMA = """
CREATE VIRTUAL TABLE IF NOT EXISTS docs USING fts5(
  path UNINDEXED, filename, content,
  tokenize='unicode61 remove_diacritics 2'
);
"""
def _sql_conn():
    con = sqlite3.connect(RAG_DB_PATH)
    con.row_factory = sqlite3.Row
    return con

def _sql_ensure():
    with _sql_conn() as con:
        con.execute(SQL_SCHEMA)

def _sql_index_folder(root: Path) -> int:
    _sql_ensure()
    cache = load_cache()
    added = 0
    with _sql_conn() as con:
        cur = con.cursor()
        for p in root.rglob("*"):
            if not p.is_file(): continue
            if any(part in SKIP_DIRS for part in p.parts): continue
            if p.suffix.lower() not in TEXT_EXT.union(DOC_EXT): continue
            sig = sig_for(p); key = str(p)
            if cache.get(key) == sig: continue
            try:
                txt = extract_text(p)
            except Exception:
                continue
            cur.execute("DELETE FROM docs WHERE path=?;", (key,))
            cur.execute("INSERT INTO docs(path, filename, content) VALUES (?,?,?);",
                        (key, p.name, txt))
            cache[key] = sig; added += 1
        con.commit()
    save_cache(cache)
    return added

def _sql_search(q: str, k:int=12) -> List[Dict]:
    _sql_ensure()
    sql = """
    SELECT path, filename,
           snippet(docs, 2, '<em>', '</em>', ' … ', 10) AS snip,
           bm25(docs) AS score, content
    FROM docs
    WHERE docs MATCH ?
    ORDER BY score ASC
    LIMIT ?;
    """
    with _sql_conn() as con:
        rows = con.execute(sql, (q, k)).fetchall()
    hits=[]
    for r in rows:
        hits.append({
            "_score": float(r["score"]),
            "_source": {"path": r["path"], "filename": r["filename"], "content": r["content"]},
            "highlight": {"content": [r["snip"]]}
        })
    return hits

# =========================================================
#                BACKEND B) Whoosh
# =========================================================
def _whoosh_ix():
    from whoosh import index
    from whoosh.fields import Schema, TEXT, ID, NUMERIC
    if not WHOOSH_INDEX_DIR.exists(): WHOOSH_INDEX_DIR.mkdir(parents=True, exist_ok=True)
    if not index.exists_in(WHOOSH_INDEX_DIR):
        schema = Schema(path=ID(unique=True, stored=True),
                        filename=TEXT(stored=True),
                        ext=ID(stored=True),
                        mtime=NUMERIC(stored=True),
                        content=TEXT(stored=True))
        return index.create_in(WHOOSH_INDEX_DIR, schema)
    return index.open_dir(WHOOSH_INDEX_DIR)

def _whoosh_index_folder(root: Path) -> int:
    from whoosh import index
    cache = load_cache()
    ix = _whoosh_ix()
    w = ix.writer(limitmb=256, procs=1, multisegment=True)
    added=0
    try:
        for p in root.rglob("*"):
            if not p.is_file(): continue
            if any(part in SKIP_DIRS for part in p.parts): continue
            if p.suffix.lower() not in TEXT_EXT.union(DOC_EXT): continue
            sig = sig_for(p); key=str(p)
            if cache.get(key) == sig: continue
            try:
                txt = extract_text(p)
            except Exception:
                continue
            w.update_document(path=key, filename=p.name, ext=p.suffix.lower(),
                              mtime=int(p.stat().st_mtime*1000), content=txt)
            cache[key]=sig; added+=1
    finally:
        w.commit()
        save_cache(cache)
    return added

def _whoosh_search(q: str, k:int=12) -> List[Dict]:
    from whoosh import scoring, highlight
    from whoosh.qparser import MultifieldParser
    ix = _whoosh_ix()
    with ix.searcher(weighting=scoring.BM25F(B=0.75, K1=1.5)) as s:
        s.formatter = highlight.HtmlFormatter(tagname="em")
        parser = MultifieldParser(["content","filename","path"], schema=ix.schema)
        rs = s.search(parser.parse(q), limit=k)
        rs.fragmenter.charlimit = None
        hits=[]
        for r in rs:
            hits.append({
                "_score": float(r.score),
                "_source": {"path": r["path"], "filename": r["filename"], "content": r["content"]},
                "highlight": {"content": [r.highlights("content")]}
            })
        return hits

# =========================================================
#                API compatible avec ton app
# =========================================================
def index_folder(root: Path) -> int:
    return _sql_index_folder(root) if RAG_MODE=="sqlite" else _whoosh_index_folder(root)

def search(question: str, k:int=12) -> List[Dict]:
    return _sql_search(question, k) if RAG_MODE=="sqlite" else _whoosh_search(question, k)

_vectorizer: Optional[SmartVectorizer] = None
def ensure_vectorizer():
    global _vectorizer
    if _vectorizer is None: _vectorizer = SmartVectorizer()
    return _vectorizer

def ask(question: str, k_bm25=12, k_final=6, use_llm=True) -> str:
    hits = search(question, k=k_bm25)
    if not hits: return "😕 Aucun document trouvé. (Index vide ?)"
    hits = ensure_vectorizer().rerank(question, hits, top_n=k_final)
    if use_llm: return answer_with_openai(question, hits)
    frags=[re.sub(r"<em>(.*?)</em>", r"**\1**", h.get("highlight",{}).get("content",[""])[0]) for h in hits]
    return ("\n\n".join(frags[:5]) or "Documents trouvés, mais aucun extrait saillant.") + \
           "\n\n**Sources**\n" + format_sources(hits)
